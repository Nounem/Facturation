#!/usr/bin/env python3
"""Génère la présentation PowerPoint éditable du projet de facturation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "presentation-facturation-salesforce.pptx"

COLORS = {
    "navy": "032D60",
    "blue": "0B5CAB",
    "sky": "1B96FF",
    "teal": "0B827C",
    "green": "2E844A",
    "orange": "B85D00",
    "red": "BA0517",
    "ink": "17202E",
    "muted": "5C6673",
    "line": "D9E1EA",
    "light": "F4F7FB",
    "white": "FFFFFF",
    "pale_blue": "EAF5FE",
    "pale_green": "E9F8F0",
    "pale_orange": "FFF3E8",
    "pale_purple": "F3EDFF",
    "purple": "6F4EB5",
}


def color(name_or_hex):
    value = COLORS.get(name_or_hex, name_or_hex).lstrip("#")
    return RGBColor.from_string(value)


def set_fill(shape, fill_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill_color)


def set_line(shape, line_color=None, width=1):
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color(line_color)
        shape.line.width = Pt(width)


def add_shape(slide, kind, x, y, w, h, fill="white", line=None, radius=True):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    font="Aptos",
    text_color="ink",
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.06,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(text_color)
    return box


def add_rich_text(slide, parts, x, y, w, h, size=18, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, bold, text_color in parts:
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color(text_color)
    return box


def add_bullets(slide, items, x, y, w, h, size=17, text_color="ink", spacing=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color(text_color)
        paragraph.space_after = Pt(spacing)
        paragraph.line_spacing = 1.05
        paragraph.text = "•  " + item
    return box


def add_header(slide, title, kicker, number):
    add_text(slide, kicker.upper(), 0.72, 0.35, 5.5, 0.3, 10, text_color="blue", bold=True)
    add_text(slide, title, 0.72, 0.72, 11.8, 0.65, 25, text_color="navy", bold=True)
    line = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 1.39, 0.72, 0.045, "sky")
    set_line(line, None)
    add_text(slide, f"{number:02d}", 12.05, 0.4, 0.55, 0.35, 10, text_color="muted", bold=True, align=PP_ALIGN.RIGHT)


def add_footer(slide):
    add_text(slide, "FACTURATION SALESFORCE  •  JUILLET 2026", 0.72, 7.14, 5.5, 0.2, 7.5, text_color="muted")
    add_text(slide, "NOUNEM / FACTURATION", 9.9, 7.14, 2.7, 0.2, 7.5, text_color="muted", align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    frame = slide.notes_slide.notes_text_frame
    frame.text = text


def add_card(slide, x, y, w, h, title, body, accent="blue", fill="white"):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, "line")
    band = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, accent)
    set_line(band, None)
    add_text(slide, title, x + 0.24, y + 0.18, w - 0.42, 0.35, 15, text_color="navy", bold=True)
    add_text(slide, body, x + 0.24, y + 0.61, w - 0.42, h - 0.76, 11.5, text_color="muted")
    return card


def add_metric(slide, x, y, w, value, label, accent, fill):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 1.28, fill, None)
    add_text(slide, value, x + 0.18, y + 0.17, w - 0.36, 0.55, 25, text_color=accent, bold=True)
    add_text(slide, label, x + 0.18, y + 0.78, w - 0.36, 0.3, 10.5, text_color="muted", bold=True)


def add_pill(slide, text, x, y, w, fill, text_color="white"):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.38, fill, None)
    add_text(slide, text, x + 0.03, y + 0.045, w - 0.06, 0.25, 9.5, text_color=text_color, bold=True, align=PP_ALIGN.CENTER)


def add_chevron(slide, x, y, w=0.42, h=0.42, fill="sky"):
    return add_shape(slide, MSO_SHAPE.CHEVRON, x, y, w, h, fill, None)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = color("navy")
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, "sky", None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 9.8, -0.3, 4.0, 2.0, "blue", None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 10.8, 5.8, 3.1, 2.0, "teal", None)
    add_pill(slide, "APPLICATION NATIVE SALESFORCE", 0.82, 0.72, 2.65, "blue")
    add_text(slide, "Facturation\nSalesforce", 0.82, 1.45, 7.8, 1.7, 34, text_color="white", bold=True)
    add_text(
        slide,
        "Du devis à la facture électronique, avec facture directe,\nrécurrence, avoirs et PDF personnalisable.",
        0.85,
        3.42,
        7.2,
        1.0,
        18,
        text_color="white",
    )
    add_text(slide, "Démonstration fonctionnelle & architecture", 0.85, 5.75, 5.5, 0.35, 12, text_color="pale_blue", bold=True)
    add_text(slide, "Juillet 2026", 0.85, 6.18, 2.0, 0.3, 10.5, text_color="pale_blue")
    add_notes(slide, "Présenter le projet comme un espace Salesforce unique qui couvre les ventes avec devis, la facture directe et la correction comptable par avoir. Durée cible de la présentation : 10 à 12 minutes, démonstration comprise.")


def add_executive_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Un socle complet, sans quitter Salesforce", "Vue d’ensemble", number)
    add_metric(slide, 0.75, 1.8, 2.8, "3", "PARCOURS DE CRÉATION", "blue", "pale_blue")
    add_metric(slide, 3.75, 1.8, 2.8, "AV", "AVOIRS COMPLETS & PARTIELS", "purple", "pale_purple")
    add_metric(slide, 6.75, 1.8, 2.8, "02:00", "FACTURATION QUOTIDIENNE", "orange", "pale_orange")
    add_metric(slide, 9.75, 1.8, 2.8, "83 %", "COUVERTURE DU PROJET", "green", "pale_green")
    add_card(slide, 0.75, 3.55, 3.75, 2.45, "Expérience fluide", "Création, recherche et actions contextualisées avec rafraîchissement immédiat et messages d’erreur exploitables.", "sky")
    add_card(slide, 4.78, 3.55, 3.75, 2.45, "Documents maîtrisés", "Numérotation à l’émission, verrouillage, PDF versionné, logo et identité visuelle par entité juridique.", "teal")
    add_card(slide, 8.55, 3.55, 3.75, 2.45, "Prêt à évoluer", "Règles inspirées de Revenue Cloud et préparation des flux Chorus Pro / Plateforme Agréée.", "orange")
    add_footer(slide)
    add_notes(slide, "Donner immédiatement les quatre preuves de valeur : trois parcours, avoirs, automatisation quotidienne et qualité technique. Préciser que 26 méthodes de test ont réussi, soit 28 exécutions Salesforce avec les méthodes de préparation.")


def add_problem_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Le problème métier", "Pourquoi ce projet", number)
    add_text(slide, "Les équipes perdent du temps lorsque la vente et la facturation vivent dans des outils séparés.", 0.78, 1.72, 6.0, 0.9, 21, text_color="navy", bold=True)
    add_bullets(
        slide,
        [
            "Ressaisies de prix, TVA et adresses client",
            "Parcours avec ou sans devis difficiles à unifier",
            "Récurrence fragile et risque de doublons",
            "Corrections comptables non traçables",
            "PDF et e-facture soumis à des règles légales",
        ],
        0.8,
        2.82,
        5.8,
        3.2,
        15,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.05, 1.8, 5.45, 4.65, "navy", None)
    add_text(slide, "OBJECTIF", 7.48, 2.25, 1.6, 0.28, 10, text_color="sky", bold=True)
    add_text(slide, "Un parcours simple\npour l’utilisateur,", 7.45, 2.85, 4.5, 1.2, 26, text_color="white", bold=True)
    add_text(slide, "contrôlé pour la comptabilité\net évolutif pour l’entreprise.", 7.45, 4.3, 4.35, 1.0, 18, text_color="pale_blue")
    add_footer(slide)
    add_notes(slide, "Insister sur le fait que certains clients ne créent jamais de devis. Le modèle doit couvrir les deux pratiques sans créer deux systèmes parallèles.")


def add_solution_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Une réponse native et configurable", "La solution", number)
    cards = [
        ("Devis", "Quote standard Salesforce", "blue", "pale_blue"),
        ("Facture directe", "Sans devis obligatoire", "teal", "pale_green"),
        ("Récurrence", "Règles et périodes configurables", "orange", "pale_orange"),
        ("Avoir", "Complet ou partiel, document dédié", "purple", "pale_purple"),
        ("PDF", "Logo, couleur, mentions, versions", "blue", "pale_blue"),
        ("E-facture", "Préparation et traçabilité du flux", "teal", "pale_green"),
    ]
    for index, (title, body, accent, fill) in enumerate(cards):
        col = index % 3
        row = index // 3
        x = 0.8 + col * 4.15
        y = 1.75 + row * 2.25
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3.75, 1.82, fill, None)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.24, y + 0.28, 0.48, 0.48, accent, None)
        add_text(slide, str(index + 1), x + 0.24, y + 0.37, 0.48, 0.18, 9, text_color="white", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.88, y + 0.27, 2.5, 0.35, 16, text_color="navy", bold=True)
        add_text(slide, body, x + 0.88, y + 0.77, 2.55, 0.58, 11.5, text_color="muted")
    add_footer(slide)
    add_notes(slide, "Le choix structurant est de conserver les devis standards Salesforce et d’utiliser un objet de facture dédié pour l’immuabilité, les snapshots client et la numérotation.")


def add_journey_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Trois chemins, un même cycle comptable", "Parcours utilisateur", number)
    sources = [
        ("Devis accepté", "Quote", "blue"),
        ("Facture directe", "Produits opportunité", "teal"),
        ("Facture planifiée", "Règle + prochaine date", "orange"),
    ]
    for i, (title, subtitle, accent) in enumerate(sources):
        y = 1.7 + i * 1.35
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, y, 3.0, 0.9, "white", "line")
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, y, 0.09, 0.9, accent, None)
        add_text(slide, title, 1.08, y + 0.14, 2.35, 0.28, 14, text_color="navy", bold=True)
        add_text(slide, subtitle, 1.08, y + 0.5, 2.35, 0.22, 9.5, text_color="muted")
        add_chevron(slide, 4.05, y + 0.24, fill="sky")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.82, 2.65, 2.55, 1.25, "navy", None)
    add_text(slide, "FACTURE\nBROUILLON", 5.12, 2.91, 1.95, 0.62, 17, text_color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_chevron(slide, 7.62, 3.05, fill="sky")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.28, 2.65, 1.75, 1.25, "blue", None)
    add_text(slide, "ÉMISSION", 8.5, 3.08, 1.32, 0.3, 14, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_chevron(slide, 10.22, 3.05, fill="sky")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.88, 2.65, 1.65, 1.25, "teal", None)
    add_text(slide, "PDF", 11.2, 3.08, 1.0, 0.3, 15, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.82, 5.18, 7.7, 0.9, "pale_purple", None)
    add_rich_text(slide, [("Facture émise  ", True, "purple"), ("→  Avoir brouillon  →  Émission AV  →  PDF d’avoir", False, "ink")], 5.12, 5.48, 7.1, 0.3, 14)
    add_footer(slide)
    add_notes(slide, "Les trois chemins convergent sur le même cycle. La correction repart ensuite de la facture émise vers un avoir : jamais de modification silencieuse du document d’origine.")


def add_workspace_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Une expérience intégrée à l’opportunité", "Lightning Web Component", number)
    # Mockup de l’application
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.75, 1.65, 8.3, 4.9, "light", "line")
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.75, 1.65, 8.3, 0.9, "navy", None)
    add_text(slide, "Devis, factures & avoirs", 1.05, 1.93, 3.7, 0.32, 16, text_color="white", bold=True)
    add_pill(slide, "Facture directe", 6.05, 1.9, 1.2, "teal")
    add_pill(slide, "Nouveau devis", 7.42, 1.9, 1.18, "sky")
    metrics = [("4", "Devis", "pale_blue"), ("2", "Acceptés", "pale_green"), ("3", "Factures", "pale_purple"), ("8 640 €", "Solde", "pale_orange")]
    for i, (value, label, fill) in enumerate(metrics):
        x = 1.02 + i * 1.9
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.82, 1.62, 0.92, fill, None)
        add_text(slide, value, x + 0.14, 2.98, 1.34, 0.3, 16, text_color="navy", bold=True)
        add_text(slide, label, x + 0.14, 3.35, 1.34, 0.2, 8.5, text_color="muted")
    add_text(slide, "FACTURES & AVOIRS", 1.03, 4.05, 2.5, 0.25, 9, text_color="blue", bold=True)
    rows = [
        ("FAC-2026-00012", "Émise", "2 400 €", "PDF"),
        ("AV-2026-00003", "Partiellement créditée", "480 €", "Ouvrir"),
        ("FAC-2026-00011", "Payée", "5 760 €", "PDF"),
    ]
    for i, (doc, status, total, action) in enumerate(rows):
        y = 4.42 + i * 0.57
        if i % 2 == 0:
            add_shape(slide, MSO_SHAPE.RECTANGLE, 0.95, y - 0.05, 7.87, 0.52, "white", None)
        add_text(slide, doc, 1.08, y + 0.04, 2.0, 0.2, 10, text_color="blue", bold=True)
        add_text(slide, status, 3.28, y + 0.04, 2.15, 0.2, 9.5, text_color="muted")
        add_text(slide, total, 5.68, y + 0.04, 1.15, 0.2, 9.5, text_color="ink", bold=True, align=PP_ALIGN.RIGHT)
        add_pill(slide, action, 7.18, y - 0.01, 1.05, "blue")
    add_card(slide, 9.45, 1.72, 3.0, 1.25, "Sans rechargement", "Les listes et compteurs se mettent à jour après chaque action.", "sky", "pale_blue")
    add_card(slide, 9.45, 3.25, 3.0, 1.25, "Erreur utile", "Le toast affiche le vrai contrôle métier ou champ manquant.", "red", "white")
    add_card(slide, 9.45, 4.78, 3.0, 1.25, "Responsive", "Le tableau devient une liste lisible sur les petits écrans.", "teal", "pale_green")
    add_footer(slide)
    add_notes(slide, "Cette maquette reprend l’organisation réelle du composant. Pendant la démo, montrer la création d’une facture directe et constater que la nouvelle ligne apparaît sans actualiser la page.")


def add_rules_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Des règles inspirées de Revenue Cloud", "Automatisation", number)
    options = [
        ("Fréquence", "Ponctuelle · Mensuelle · Trimestrielle · Semestrielle · Annuelle"),
        ("Moment", "À échoir ou à terme échu"),
        ("Frontière", "Période civile ou date anniversaire"),
        ("Paiement", "Jour de facturation et délai configurable"),
    ]
    for i, (title, body) in enumerate(options):
        add_card(slide, 0.78, 1.7 + i * 1.13, 5.05, 0.92, title, body, ["blue", "teal", "orange", "purple"][i], "white")
    add_text(slide, "Exemple — mensuelle à terme échu", 6.45, 1.75, 5.7, 0.42, 17, text_color="navy", bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.72, 3.17, 5.15, 0.05, "line", None)
    timeline = [
        ("1 juin", "Début période", "blue"),
        ("30 juin", "Fin prestation", "teal"),
        ("1 juillet", "Facture générée", "orange"),
        ("1 août", "Prochaine date", "purple"),
    ]
    for i, (date, label, accent) in enumerate(timeline):
        x = 6.55 + i * 1.72
        add_shape(slide, MSO_SHAPE.OVAL, x, 2.91, 0.54, 0.54, accent, None)
        add_text(slide, date, x - 0.18, 3.62, 0.9, 0.24, 10, text_color="navy", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x - 0.34, 4.02, 1.25, 0.52, 9.5, text_color="muted", align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.45, 5.12, 5.8, 0.88, "pale_green", None)
    add_text(slide, "Clé unique opportunité + échéance = protection anti-doublon", 6.78, 5.4, 5.12, 0.3, 12, text_color="green", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "Expliquer l’exemple calendrier : le 1er juillet, le système facture le mois de juin puis positionne la prochaine date au 1er août. Le batch quotidien ne recrée pas la même échéance.")


def add_legal_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Une facture stable, une correction traçable", "Conformité métier", number)
    add_text(slide, "FACTURE", 0.9, 1.72, 2.3, 0.32, 12, text_color="blue", bold=True)
    steps = [
        ("1", "Snapshot client", "Adresse et identité figées"),
        ("2", "Validation", "Mentions et lignes contrôlées"),
        ("3", "Numérotation", "Attribuée uniquement à l’émission"),
        ("4", "Verrouillage", "Document comptable immuable"),
    ]
    for i, (num, title, body) in enumerate(steps):
        y = 2.2 + i * 0.9
        add_shape(slide, MSO_SHAPE.OVAL, 0.92, y, 0.45, 0.45, "blue", None)
        add_text(slide, num, 0.92, y + 0.11, 0.45, 0.16, 8.5, text_color="white", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.62, y - 0.02, 2.2, 0.27, 13, text_color="navy", bold=True)
        add_text(slide, body, 1.62, y + 0.34, 3.5, 0.25, 10.5, text_color="muted")
    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.15, 1.75, 0.05, 4.65, "line", None)
    add_text(slide, "AVOIR", 6.65, 1.72, 2.3, 0.32, 12, text_color="purple", bold=True)
    flow = [
        ("Facture émise", "navy"),
        ("Avoir brouillon", "purple"),
        ("Contrôle du restant", "orange"),
        ("Émission AV", "teal"),
    ]
    for i, (label, accent) in enumerate(flow):
        x = 6.65 + (i % 2) * 2.85
        y = 2.3 + (i // 2) * 1.45
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.25, 0.85, accent, None)
        add_text(slide, label, x + 0.2, y + 0.26, 1.85, 0.28, 12, text_color="white", bold=True, align=PP_ALIGN.CENTER)
        if i in (0, 2):
            add_chevron(slide, x + 2.38, y + 0.21, fill="sky")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.65, 5.32, 5.1, 0.78, "pale_purple", None)
    add_text(slide, "Solde et statut de la facture recalculés automatiquement", 6.95, 5.58, 4.5, 0.25, 11.5, text_color="purple", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "Le message clé : une facture émise ne se corrige pas directement. L’avoir conserve le lien, le motif, sa propre séquence et son propre PDF. Un verrou sur la facture d’origine empêche le sur-crédit concurrent.")


def add_pdf_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Un PDF professionnel et personnalisable", "Document client", number)
    # Feuille PDF simplifiée
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.95, 1.62, 5.05, 4.95, "white", "line")
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.2, 1.88, 4.55, 0.08, "blue", None)
    add_text(slide, "VOTRE LOGO", 1.25, 2.18, 1.55, 0.28, 10, text_color="muted", bold=True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.15, 2.12, 1.35, 0.8, "blue", None)
    add_text(slide, "FACTURE", 4.28, 2.36, 1.08, 0.25, 14, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.25, 3.22, 2.05, 0.82, "light", None)
    add_text(slide, "FACTURÉ À\nClient Exemple", 1.45, 3.42, 1.65, 0.4, 9.5, text_color="navy", bold=True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.52, 3.22, 1.98, 0.82, "light", None)
    add_text(slide, "DATE  •  ÉCHÉANCE\nN° FAC-2026-00012", 3.67, 3.42, 1.65, 0.4, 8.5, text_color="navy", bold=True)
    for i in range(4):
        y = 4.35 + i * 0.38
        add_shape(slide, MSO_SHAPE.RECTANGLE, 1.25, y, 4.25, 0.04 if i == 0 else 0.02, "blue" if i == 0 else "line", None)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.78, 5.9, 1.72, 0.42, "blue", None)
    add_text(slide, "TOTAL TTC   2 400 €", 3.91, 6.03, 1.46, 0.16, 8.5, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    callouts = [
        ("Logo Salesforce Files", "Une identité par entité juridique", "blue", "pale_blue"),
        ("Couleur #RRGGBB", "Déclinée sur titres, tableau et totaux", "teal", "pale_green"),
        ("Mentions françaises", "SIREN, SIRET, TVA, pénalités, paiement", "orange", "pale_orange"),
        ("Versionnement", "Chaque régénération conserve l’historique", "purple", "pale_purple"),
    ]
    for i, (title, body, accent, fill) in enumerate(callouts):
        x = 6.6 + (i % 2) * 2.95
        y = 1.8 + (i // 2) * 2.15
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.62, 1.62, fill, None)
        add_text(slide, title, x + 0.2, y + 0.2, 2.2, 0.35, 13, text_color=accent, bold=True)
        add_text(slide, body, x + 0.2, y + 0.73, 2.2, 0.58, 10.2, text_color="muted")
    add_pill(slide, "FACTURE", 6.65, 6.05, 1.15, "blue")
    add_pill(slide, "AVOIR", 7.98, 6.05, 1.15, "purple")
    add_text(slide, "Deux mises en page adaptées au document", 9.35, 6.1, 2.75, 0.24, 10, text_color="muted")
    add_footer(slide)
    add_notes(slide, "Le PDF est généré côté serveur, attaché à la facture et versionné. L’avoir affiche en plus la facture d’origine, le motif et le montant porté au crédit du client.")


def add_einvoice_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Facturation électronique : prête à raccorder", "Chorus Pro & Plateforme Agréée", number)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 2.35, 2.3, 1.05, "navy", None)
    add_text(slide, "FACTURE ÉMISE", 1.18, 2.75, 1.72, 0.25, 14, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_chevron(slide, 3.48, 2.66, fill="sky")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.18, 2.35, 2.6, 1.05, "blue", None)
    add_text(slide, "PRÉPARER\nLA E-FACTURE", 4.58, 2.58, 1.8, 0.5, 14, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_chevron(slide, 7.02, 2.66, fill="sky")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.72, 1.68, 2.22, 1.15, "teal", None)
    add_text(slide, "CLIENT PUBLIC\nChorus Pro", 8.0, 1.98, 1.66, 0.48, 12.5, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.72, 3.16, 2.22, 1.15, "purple", None)
    add_text(slide, "CLIENT PRIVÉ\nPlateforme Agréée", 7.95, 3.45, 1.76, 0.48, 12, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.35, 2.35, 2.05, 1.05, "pale_green", None)
    add_text(slide, "TRANSMISSION\nTRAÇABLE", 10.61, 2.6, 1.52, 0.45, 12.5, text_color="green", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "1er SEPTEMBRE 2026", 0.9, 5.05, 2.75, 0.3, 13, text_color="green", bold=True)
    add_text(slide, "Réception pour tous · émission grandes entreprises et ETI", 0.9, 5.5, 5.4, 0.3, 11.5, text_color="muted")
    add_text(slide, "1er SEPTEMBRE 2027", 7.05, 5.05, 2.75, 0.3, 13, text_color="orange", bold=True)
    add_text(slide, "Émission obligatoire pour les PME et micro-entreprises", 7.05, 5.5, 5.15, 0.3, 11.5, text_color="muted")
    add_footer(slide)
    add_notes(slide, "Être transparent : la préparation, le routage et la traçabilité existent. L’envoi externe reste volontairement désactivé tant que le raccordement PISTE et la conformité ne sont pas certifiés. Un PDF ordinaire ne suffit pas à constituer une facture électronique conforme.")


def add_architecture_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Une architecture simple à faire évoluer", "Vue technique", number)
    layers = [
        ("EXPÉRIENCE", "billingWorkspace · brandingConfigurator", "blue", "pale_blue"),
        ("FAÇADE", "BillingWorkspaceController · BillingSecurity", "purple", "pale_purple"),
        ("MÉTIER", "InvoiceService · Validation · Calcul · Numérotation", "teal", "pale_green"),
        ("AUTOMATISATION", "Scheduler · Batch · Règles · Idempotence", "orange", "pale_orange"),
        ("DONNÉES & DOCUMENTS", "Invoice · Lines · Payments · Files · Transmissions", "navy", "light"),
    ]
    for i, (title, body, accent, fill) in enumerate(layers):
        y = 1.66 + i * 1.0
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.05 + i * 0.26, y, 11.0 - i * 0.52, 0.75, fill, None)
        add_pill(slide, title, 1.3 + i * 0.26, y + 0.18, 1.72, accent)
        add_text(slide, body, 3.35 + i * 0.22, y + 0.22, 7.6 - i * 0.4, 0.25, 12, text_color="ink", bold=True)
    add_footer(slide)
    add_notes(slide, "La séparation expérience, façade, services et données permet de faire évoluer le LWC ou le connecteur externe sans réécrire les calculs. Toutes les classes respectent le partage Salesforce.")


def add_quality_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Qualité, sécurité et exploitation", "Preuves de livraison", number)
    add_shape(slide, MSO_SHAPE.OVAL, 0.95, 1.82, 2.35, 2.35, "pale_green", None)
    add_text(slide, "83 %", 1.2, 2.42, 1.85, 0.62, 30, text_color="green", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "couverture projet", 1.3, 3.2, 1.65, 0.28, 10, text_color="green", bold=True, align=PP_ALIGN.CENTER)
    add_metric(slide, 3.72, 1.9, 2.25, "100 %", "TESTS RÉUSSIS", "blue", "pale_blue")
    add_metric(slide, 6.2, 1.9, 2.25, "26", "MÉTHODES DE TEST", "purple", "pale_purple")
    add_metric(slide, 8.68, 1.9, 2.25, "0", "ÉCHEC", "green", "pale_green")
    add_card(slide, 0.95, 4.65, 3.5, 1.4, "Sécurité", "Partage Salesforce, contrôles CRUD/FLS et Permission Set dédié.", "blue")
    add_card(slide, 4.76, 4.65, 3.5, 1.4, "Résilience", "Transactions, verrouillage, anti-doublon et erreurs isolées par opportunité.", "teal")
    add_card(slide, 8.57, 4.65, 3.5, 1.4, "Supervision", "Scheduler vérifiable, erreurs de batch et transmissions historisées.", "orange")
    add_text(slide, "Run Salesforce : 707KB000027XkfR", 0.98, 6.35, 4.2, 0.23, 9.5, text_color="muted")
    add_footer(slide)
    add_notes(slide, "La couverture de 83 % concerne le périmètre de classes du projet exécuté par cette suite. La couverture globale de l’organisation contient d’autres développements historiques et ne doit pas être confondue avec cette mesure.")


def add_delivery_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Une livraison traçable par GitHub", "Déploiement", number)
    flow = [
        ("COMMIT", "Version source", "navy"),
        ("VALIDATE", "Sans écriture", "blue"),
        ("TESTS", "Run spécifié", "purple"),
        ("APPROBATION", "Mot DEPLOY", "orange"),
        ("SALESFORCE", "Déploiement", "teal"),
    ]
    for i, (title, body, accent) in enumerate(flow):
        x = 0.75 + i * 2.48
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.15, 1.9, 1.25, accent, None)
        add_text(slide, title, x + 0.15, 2.48, 1.6, 0.28, 12, text_color="white", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.15, 2.92, 1.6, 0.2, 8.5, text_color="white", align=PP_ALIGN.CENTER)
        if i < 4:
            add_chevron(slide, x + 2.02, 2.56, fill="sky")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.55, 4.35, 10.2, 1.35, "light", None)
    add_text(slide, "Version livrée", 1.9, 4.66, 1.7, 0.3, 12, text_color="muted", bold=True)
    add_text(slide, "Commit 1b184f4", 3.62, 4.63, 2.15, 0.34, 15, text_color="blue", bold=True)
    link = add_text(slide, "github.com/Nounem/Facturation", 6.02, 4.65, 4.6, 0.3, 13, text_color="blue", bold=True)
    link.text_frame.paragraphs[0].runs[0].hyperlink.address = "https://github.com/Nounem/Facturation"
    add_text(slide, "Déploiement Salesforce : 0AfKB00000DphNT0AZ", 3.62, 5.12, 5.0, 0.25, 10, text_color="muted")
    add_footer(slide)
    add_notes(slide, "GitHub est la source de vérité. Le workflow valide les métadonnées et les tests avant toute écriture ; le mot DEPLOY et la protection d’environnement apportent une approbation explicite.")


def add_demo_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Démonstration proposée — 7 minutes", "Scénario", number)
    steps = [
        ("01", "Ouvrir l’opportunité", "Compte + produits prêts"),
        ("02", "Créer une facture directe", "La ligne apparaît sans recharger"),
        ("03", "Émettre", "Contrôles, numéro et verrouillage"),
        ("04", "Générer le PDF", "Logo, couleur et fichier versionné"),
        ("05", "Créer un avoir partiel", "Solde et statut recalculés"),
        ("06", "Montrer la règle", "Prochaine échéance et anti-doublon"),
    ]
    for i, (num, title, body) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = 0.9 + col * 6.05
        y = 1.7 + row * 1.55
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 5.45, 1.15, "white", "line")
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.22, y + 0.25, 0.62, 0.62, "blue" if col == 0 else "teal", None)
        add_text(slide, num, x + 0.22, y + 0.43, 0.62, 0.18, 8.5, text_color="white", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 1.08, y + 0.22, 3.9, 0.3, 14, text_color="navy", bold=True)
        add_text(slide, body, x + 1.08, y + 0.64, 3.95, 0.27, 10.5, text_color="muted")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.75, 6.32, 7.8, 0.45, "pale_orange", None)
    add_text(slide, "Préparer les données en amont et ne pas dépendre d’un appel Chorus réel pendant la démo.", 3.05, 6.44, 7.2, 0.18, 9.5, text_color="orange", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "Préparer une opportunité avec deux produits, une entité juridique complète et les séquences de l’année. Le parcours doit rester court et visible : création, émission, PDF, avoir puis règle.")


def add_roadmap_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Le socle est livré, la suite est maîtrisée", "Prochaines étapes", number)
    add_text(slide, "COURT TERME", 0.92, 1.75, 2.2, 0.3, 11, text_color="green", bold=True)
    add_bullets(slide, ["Renseigner les vraies mentions juridiques", "Recette avec la comptabilité", "Valider le logo et le modèle PDF", "Finaliser le raccordement PISTE / PA"], 0.92, 2.25, 5.2, 2.6, 15)
    add_text(slide, "ÉVOLUTIONS", 6.72, 1.75, 2.2, 0.3, 11, text_color="purple", bold=True)
    add_bullets(slide, ["Contrats et échéanciers par ligne", "Relances et prélèvements", "Allocation multi-paiements", "Portail client et suivi du règlement"], 6.72, 2.25, 5.2, 2.6, 15)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 5.3, 10.3, 0.95, "navy", None)
    add_text(slide, "Un système simple aujourd’hui, extensible demain vers un modèle Revenue Cloud.", 1.92, 5.61, 9.45, 0.3, 17, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, "Conclure sur un socle utilisable immédiatement. Ne pas présenter l’envoi électronique comme certifié ; c’est le principal chantier externe restant.")


def add_final_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = color("navy")
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, "sky", None)
    add_text(slide, "Merci", 0.92, 1.65, 5.8, 0.85, 36, text_color="white", bold=True)
    add_text(slide, "Questions ?", 0.95, 2.75, 4.2, 0.55, 23, text_color="pale_blue", bold=True)
    add_text(slide, "Documentation technique, guide utilisateur et workflow\ndisponibles dans le dépôt GitHub.", 0.95, 4.05, 6.7, 0.8, 16, text_color="white")
    link = add_text(slide, "github.com/Nounem/Facturation", 0.95, 5.28, 5.5, 0.35, 13, text_color="sky", bold=True)
    link.text_frame.paragraphs[0].runs[0].hyperlink.address = "https://github.com/Nounem/Facturation"
    add_shape(slide, MSO_SHAPE.OVAL, 9.3, 1.45, 2.55, 2.55, "blue", None)
    add_text(slide, "SF", 9.75, 2.12, 1.65, 0.75, 34, text_color="white", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "FACTURATION", 9.75, 4.52, 1.65, 0.28, 11, text_color="sky", bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "Garder ouvertes la page Opportunité, une facture émise et la page GitHub Actions pour répondre aux questions avec des preuves concrètes.")


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Facturation Salesforce"
    prs.core_properties.subject = "Présentation fonctionnelle et technique"
    prs.core_properties.author = "Nounem"
    prs.core_properties.keywords = "Salesforce, facturation, devis, avoir, PDF, Chorus Pro"
    prs.core_properties.comments = "Généré depuis scripts/generate_presentation.py"

    add_title_slide(prs)
    builders = [
        add_executive_slide,
        add_problem_slide,
        add_solution_slide,
        add_journey_slide,
        add_workspace_slide,
        add_rules_slide,
        add_legal_slide,
        add_pdf_slide,
        add_einvoice_slide,
        add_architecture_slide,
        add_quality_slide,
        add_delivery_slide,
        add_demo_slide,
        add_roadmap_slide,
    ]
    for number, builder in enumerate(builders, start=2):
        builder(prs, number)
    add_final_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return prs


if __name__ == "__main__":
    presentation = build_presentation()
    print(f"Présentation générée : {OUTPUT} ({len(presentation.slides)} diapositives)")
